import os
import pdfplumber
from docx import Document
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract

DOSSIER_DOCUMENTS = "documents"

def lire_pdf(chemin):
    texte = ""
    with pdfplumber.open(chemin) as pdf:
        for page in pdf.pages:
            contenu = page.extract_text()
            if contenu:
                texte += contenu + "\n"
    return texte

def lire_pdf_scanné(chemin):
    texte = ""
    images = convert_from_path(chemin)
    for image in images:
        contenu = pytesseract.image_to_string(image, lang="fra")
        texte += contenu + "\n"
    return texte

def lire_word(chemin):
    doc = Document(chemin)
    texte = ""
    for paragraphe in doc.paragraphs:
        texte += paragraphe.text + "\n"
    return texte

def lire_powerpoint(chemin):
    prs = Presentation(chemin)
    texte = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texte += shape.text + "\n"
    return texte

def charger_documents():
    documents = []
    for fichier in os.listdir(DOSSIER_DOCUMENTS):
        chemin = os.path.join(DOSSIER_DOCUMENTS, fichier)
        if fichier.endswith(".pdf"):
            print(f"Lecture PDF : {fichier}")
            texte = lire_pdf(chemin)
            if len(texte.strip()) < 100:
                print(f"  → PDF scanné détecté, OCR en cours...")
                texte = lire_pdf_scanné(chemin)
            documents.append({"source": fichier, "texte": texte})
        elif fichier.endswith(".docx"):
            print(f"Lecture Word : {fichier}")
            texte = lire_word(chemin)
            documents.append({"source": fichier, "texte": texte})
        elif fichier.endswith(".pptx"):
            print(f"Lecture PowerPoint : {fichier}")
            texte = lire_powerpoint(chemin)
            documents.append({"source": fichier, "texte": texte})
    return documents

if __name__ == "__main__":
    os.makedirs(DOSSIER_DOCUMENTS, exist_ok=True)
    docs = charger_documents()
    print(f"\n{len(docs)} document(s) chargé(s).")
    for doc in docs:
        print(f"- {doc['source']} : {len(doc['texte'])} caractères extraits")